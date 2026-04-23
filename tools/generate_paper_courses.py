from __future__ import annotations

import json
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(r"g:/pythonProject/deliverables/paper-to-course")


COURSES = [
    {
        "name": "2024_bnns_sf_cooling_patch",
        "title": "High Thermal Conductivity and Radiative Cooling Designed Boron Nitride Nanosheets/Silk Fibroin Films for Personal Thermal Management",
        "subtitle": "ACS Applied Materials & Interfaces 2024",
        "authors": "Juncheng Xia et al.",
        "doi": "10.1021/acsami.3c16602",
        "url": "https://pubs.acs.org/doi/10.1021/acsami.3c16602",
        "summary": "BNNS and silk fibroin are combined into an aligned, reusable cooling film that balances heat spreading, radiative cooling and skin compatibility.",
        "source_note": "ACS abstract, article page, figures 1-4 and conclusion.",
        "domain": "personal thermal management / radiative cooling / biobased films",
        "timeline": [
            (2012, "Silk thermal transport", "Silk becomes a serious thermal-material candidate."),
            (2017, "BNNS composite films", "Oriented BN starts to dominate in-plane heat flow."),
            (2020, "Wearable passive cooling", "Textiles and films converge for body cooling."),
            (2024, "BNNS/SF film", "The paper integrates alignment, reflectance and biocompatibility."),
        ],
        "compare": [
            ["Pure SF", "Biobased and flexible", "Thermal transport is too low"],
            ["Graphene textiles", "Good conduction", "Wearability and outdoor cooling tradeoffs"],
            ["BNNS polymer films", "Strong anisotropic heat spreading", "Less skin-friendly matrix"],
            ["BNNS/SF film", "Aligned, reusable, biocompatible", "High filler loading may reduce softness"],
        ],
        "cards": [
            ["Why it matters", "Cooling is energy intensive, so passive skin-contact cooling is a practical target."],
            ["Main bottleneck", "Pure silk is biocompatible but not conductive enough for rapid heat dissipation."],
            ["Core idea", "Use BNNS as a directional heat spreader and silk fibroin as a durable matrix."],
        ],
        "steps": [
            ["1", "Prepare BNNS", "Exfoliate h-BN into nanosheets without losing the platelet structure."],
            ["2", "Mix with SF", "Combine with silk fibroin before vacuum filtration."],
            ["3", "Build alignment", "Layer-by-layer stacking forms a directional thermal path."],
            ["4", "Validate", "Measure thermal conductivity, reflectance, emissivity and biosafety."],
        ],
        "formula": "λ = α × ρ × Cp; TCE = (λcomposite - λSF) / λSF × 100%",
        "formula_lines": [
            ["λ", "effective thermal conductivity"],
            ["α", "thermal diffusivity"],
            ["ρ", "density"],
            ["Cp", "specific heat capacity"],
        ],
        "stats": [
            ["27.58 W m−1 K−1", "in-plane thermal conductivity"],
            ["1.77 W m−1 K−1", "out-of-plane thermal conductivity"],
            ["0.89", "solar reflectance"],
            ["0.84", "LWIR emissivity"],
        ],
        "bars": [
            ["Pure SF", 8, "baseline"],
            ["BNNS/SF 50 wt%", 100, "best conduction"],
            ["Indoor cooling", 92, "largest effect"],
            ["Outdoor cooling", 85, "still strong"],
        ],
        "limitations": [
            "Very high filler loading may reduce softness.",
            "Large-area alignment is harder to scale.",
            "Wear, sweat and washing need broader field testing.",
        ],
        "future": [
            "Optimize interfacial chemistry to reduce phonon scattering.",
            "Develop roll-to-roll or other scalable alignment routes.",
            "Test laundering and fatigue under realistic clothing conditions.",
            "Co-design with textile architectures for all-day comfort.",
        ],
        "quiz": {
            "question": "Which feature is the main reason BNNS helps this film?",
            "answer": "anisotropic_heat_spreading",
            "options": [
                ["It increases anisotropic heat spreading", "anisotropic_heat_spreading"],
                ["It only adds color", "color"],
                ["It removes all interface resistance", "resistance"],
            ],
        },
    },
    {
        "name": "2020_silk_bn_nanofibers",
        "title": "Electrospun Silk-Boron Nitride Nanofibers with Tunable Structure and Properties",
        "subtitle": "Polymers 2020",
        "authors": "Ye Xue and Xiao Hu",
        "doi": "10.3390/polym12051093",
        "url": "https://www.mdpi.com/2073-4360/12/5/1093",
        "summary": "h-BN changes the morphology, secondary structure and thermal stability of electrospun silk fibers, making them more tunable for biofriendly composite use.",
        "source_note": "MDPI abstract, article page, figures 1-5 and conclusion.",
        "domain": "electrospun biocomposites / silk fibroin / h-BN fibers",
        "timeline": [
            (2003, "Electrospinning matures", "Fiber morphology becomes controllable by process tuning."),
            (2010, "Silk biomaterials", "Silk expands into scaffolds, dressings and devices."),
            (2015, "BN thermal fillers", "BN becomes a major electrically insulating conductor filler."),
            (2020, "Silk-BN fibers", "The paper maps how BN changes fiber structure and stability."),
        ],
        "compare": [
            ["Pure SF nanofibers", "Smooth and biofriendly", "Limited thermal robustness"],
            ["Low-BN BNSF", "Better alignment", "Modest reinforcement"],
            ["High-BN BNSF", "Rougher but more coupled", "Excess loading may hurt texture"],
            ["This work", "Tunable electrospun biohybrid", "No direct conductivity target reported"],
        ],
        "cards": [
            ["Why it matters", "Electrospun silk can be made into porous mats, but its thermal and structural stability still needs help."],
            ["Main bottleneck", "A filler must reinforce the protein matrix without destroying processability."],
            ["Core idea", "Use h-BN nanosheets to tune fiber morphology, crosslinking and protein secondary structure."],
        ],
        "steps": [
            ["1", "Dissolve silk", "Prepare a formic-acid based spinning dope."],
            ["2", "Add BN", "Mix h-BN at different loading ratios."],
            ["3", "Electrospin", "Form composite fibers with tunable diameter and packing."],
            ["4", "Characterize", "Use SEM, FTIR, TGA and TM-DSC to map structure-property changes."],
        ],
        "formula": "Thermal stability is tracked by Tw, Tg and Td shifts after BN addition.",
        "formula_lines": [
            ["Tw", "bound-water evaporation feature"],
            ["Tg", "glass transition temperature"],
            ["Td", "thermal degradation feature"],
            ["Diameter", "morphology marker for electrospun mats"],
        ],
        "stats": [
            ["800 nm", "approx. diameter of pure SF fibers"],
            ["400-750 nm", "diameter range after BN addition"],
            ["120 °C", "pure SF onset Tg"],
            ["96 °C", "lowest Tg at 40% BN"],
        ],
        "bars": [
            ["Pure SF Tg", 100, "highest"],
            ["40% BNSF Tg", 80, "lower"],
            ["Residual mass", 92, "improves"],
            ["Morphology control", 88, "strong"],
        ],
        "limitations": [
            "The paper emphasizes structure tuning more than conductivity benchmarking.",
            "Higher loading can roughen the fibers.",
            "Biomedical validation remains prospective.",
        ],
        "future": [
            "Measure conductivity and anisotropy explicitly in aligned mats.",
            "Link morphology tuning to cell response in wet environments.",
            "Explore post-spinning annealing or stretching to amplify ordering.",
            "Combine BN with additional benign fillers for multifunctionality.",
        ],
        "quiz": {
            "question": "What was the strongest reported effect of BN on the silk nanofibers?",
            "answer": "thermal_stability_and_structure",
            "options": [
                ["Thermal stability and structure tuning", "thermal_stability_and_structure"],
                ["Magnetic switching", "magnetic"],
                ["Superconductivity", "superconductivity"],
            ],
        },
    },
    {
        "name": "2025_cd_bnoh_silma_patch",
        "title": "High Thermal Conductivity and Biocompatibility Silk Fibroin Hydrogels Filled with CDs@BNOH for Body Surface Cooling",
        "subtitle": "ACS Applied Materials & Interfaces 2025",
        "authors": "Shuming Liu et al.",
        "doi": "10.1021/acsami.5c09260",
        "url": "https://pubs.acs.org/doi/10.1021/acsami.5c09260",
        "summary": "Carbon-dot-modified hydroxylated BN is embedded into a silk hydrogel patch to balance thermal conduction, UV handling, biodegradability and skin-contact cooling.",
        "source_note": "ACS abstract, article page, figures 1-5 and conclusion.",
        "domain": "body surface cooling / hydrogel patch / carbon-dot-modified BN",
        "timeline": [
            (2020, "Wearable passive cooling", "Hydrogels and textiles begin to merge with cooling goals."),
            (2022, "Biocompatible gels", "Skin-contact materials become a major design target."),
            (2024, "BNNS/SF cooling films", "Film-type alignment and radiative cooling are proven."),
            (2025, "CDs@BNOH/SilMA", "This work adds surface engineering and degradable hydrogel integration."),
        ],
        "compare": [
            ["Commercial H2O@PAAS patch", "Simple and available", "Depends heavily on evaporation"],
            ["Pure SilMA hydrogel", "Soft and biobased", "Thermal conductivity is modest"],
            ["BNOH/SilMA", "Better conduction than pure gel", "Less interface engineering"],
            ["CDs@BNOH/SilMA", "Best balance here", "Long-term wear still needs more data"],
        ],
        "cards": [
            ["Why it matters", "Wearable cooling needs a matrix that is biocompatible, degradable and thermally useful."],
            ["Main bottleneck", "Many conductive fillers are good conductors but poor skin-contact materials without surface engineering."],
            ["Core idea", "Hydroxylated BN is further modified by carbon quantum dots to improve phonon transport and UV response."],
        ],
        "steps": [
            ["1", "Prepare CDs", "Make carbon quantum dots from urea/citrate precursors."],
            ["2", "Build CDs@BNOH", "Functionalize hydroxylated BN with CDs to tune the interface."],
            ["3", "Make SilMA", "Methacrylate-modify silk fibroin for UV-curable gelation."],
            ["4", "Infuse patch", "Load the filler into cotton-supported gel and cure into a wearable patch."],
        ],
        "formula": "Proper CDs modification improves phonon transport; too much coating disrupts the BN network.",
        "formula_lines": [
            ["fwhm", "Raman linewidth used as a phonon-lifetime indicator"],
            ["D%", "degradation fraction of the hydrogel"],
            ["ΔT", "surface cooling effect relative to blank"],
            ["λ", "patch thermal conductivity"],
        ],
        "stats": [
            ["0.77 W m−1 K−1", "patch thermal conductivity"],
            ["15.38%", "overall conductivity increase"],
            ["6.6 °C", "cooling improvement vs. control"],
            ["10 °C+", "indoor scenario peak cooling effect"],
        ],
        "bars": [
            ["BNOH/SilMA", 86, "baseline conductive gel"],
            ["CDs@BNOH/SilMA", 100, "best"],
            ["Commercial patch", 55, "evaporation-dependent"],
            ["UV shielding", 90, "good"],
        ],
        "limitations": [
            "Conductivity is still below the best film-type BNNS systems.",
            "Outdoor performance depends on contact and ambient conditions.",
            "Long-term wear tests are still limited.",
        ],
        "future": [
            "Quantify comfort under sweat, motion and repeated bending.",
            "Optimize CDs loading to keep the BN phonon network intact.",
            "Explore device integration for wearable electronics.",
            "Test other biodegradable matrices and patch geometries.",
        ],
        "quiz": {
            "question": "Why do the carbon dots matter in the filler?",
            "answer": "interface_tuning",
            "options": [
                ["They tune the BN interface and help transport", "interface_tuning"],
                ["They only make the patch heavier", "heavier"],
                ["They remove all degradation", "remove_degradation"],
            ],
        },
    },
]


CSS = ":root{--bg:#faf7f2;--panel:#fffaf4;--text:#2d2620;--muted:#6b6057;--accent:#d94f30;--line:#eadfce;--shadow:0 10px 30px rgba(45,38,32,.08)}*{box-sizing:border-box}body{margin:0;font-family:Arial,Helvetica,sans-serif;background:linear-gradient(180deg,#fcfbf8 0%,var(--bg) 100%);color:var(--text);line-height:1.6}a{color:var(--accent);text-decoration:none}.wrap{max-width:1160px;margin:0 auto;padding:0 20px 40px}.hero{padding:40px 0 22px;border-bottom:1px solid var(--line)}.kicker{text-transform:uppercase;letter-spacing:.18em;color:var(--accent);font-size:12px;font-weight:700}h1{font-size:clamp(30px,4vw,52px);line-height:1.05;margin:10px 0 12px}.lead{max-width:860px;color:var(--muted);font-size:17px}.meta{display:flex;flex-wrap:wrap;gap:10px;margin-top:18px}.pill{background:rgba(217,79,48,.09);color:#8a4d36;border:1px solid rgba(217,79,48,.18);border-radius:999px;padding:7px 12px;font-size:13px}.nav{position:sticky;top:0;z-index:20;background:rgba(250,247,242,.9);backdrop-filter:blur(8px);border-bottom:1px solid var(--line)}.nav-inner{display:flex;flex-wrap:wrap;gap:10px;padding:12px 0}.nav a{padding:8px 12px;border-radius:10px;color:var(--text);font-size:14px;border:1px solid transparent}.nav a.active,.nav a:hover{background:#fff;border-color:var(--line);box-shadow:var(--shadow)}.module{background:rgba(255,255,255,.86);border:1px solid var(--line);border-radius:20px;box-shadow:var(--shadow);padding:24px;margin-top:18px}.module h2{margin:0 0 8px;font-size:26px}.module p{color:var(--muted)}.grid{display:grid;gap:18px}.cards{grid-template-columns:repeat(auto-fit,minmax(220px,1fr))}.two-col{grid-template-columns:1.1fr .9fr}.card{background:var(--panel);border:1px solid var(--line);border-radius:16px;padding:18px}.card h3{margin:0 0 8px;font-size:18px}.table-wrap{overflow-x:auto;border-radius:14px;border:1px solid var(--line)}table{width:100%;border-collapse:collapse;background:#fff}th,td{padding:12px 14px;border-bottom:1px solid #eee7de;text-align:left;vertical-align:top}th{background:#fcf6ef}.timeline{position:relative;margin:14px 0;padding-left:24px}.timeline:before{content:'';position:absolute;left:9px;top:4px;bottom:4px;width:2px;background:rgba(217,79,48,.22)}.timeline-item{position:relative;margin-bottom:16px}.timeline-item:before{content:'';position:absolute;left:-24px;top:6px;width:18px;height:18px;border-radius:50%;background:var(--accent)}.timeline-year{font-weight:800;color:var(--accent)}.quote{border-left:4px solid var(--accent);padding:12px 14px;background:#fff6f0;border-radius:0 12px 12px 0}.quiz button{border:1px solid var(--line);background:#fff;color:var(--text);border-radius:12px;padding:10px 14px;cursor:pointer;margin:4px 6px 4px 0}.quiz button.correct{background:#e9f7ef;border-color:#9fdfb4}.quiz button.wrong{background:#fdeaea;border-color:#f1b4b4}.footer{margin-top:24px;padding:20px 0 12px;color:var(--muted);border-top:1px solid var(--line);font-size:14px}.badge{display:inline-block;padding:4px 10px;border-radius:999px;background:#fff;border:1px solid var(--line);margin-right:6px}@media(max-width:900px){.two-col{grid-template-columns:1fr}.hero{padding-top:26px}}"


JS = "window.addEventListener('DOMContentLoaded',()=>{const links=[...document.querySelectorAll('.nav a')];const sections=links.map(link=>document.querySelector(link.getAttribute('href'))).filter(Boolean);const observer=new IntersectionObserver(entries=>{const visible=entries.filter(entry=>entry.isIntersecting).sort((a,b)=>b.intersectionRatio-a.intersectionRatio)[0];if(!visible)return;links.forEach(link=>link.classList.toggle('active',link.getAttribute('href')==='#'+visible.target.id))},{threshold:[.2,.4,.6]});sections.forEach(sec=>observer.observe(sec));document.querySelectorAll('[data-quiz]').forEach(box=>{const answer=box.dataset.answer;box.querySelectorAll('button').forEach(btn=>btn.addEventListener('click',()=>{box.querySelectorAll('button').forEach(b=>b.classList.remove('correct','wrong'));btn.classList.add(btn.dataset.choice===answer?'correct':'wrong')}))})});"


BASE_HTML = """<!doctype html><html lang='zh-CN'><head><meta charset='utf-8' /><meta name='viewport' content='width=device-width, initial-scale=1' /><title>{title}</title><link rel='stylesheet' href='styles.css' /></head><body><div class='nav'><div class='wrap nav-inner'>{nav}</div></div><div class='wrap'><header class='hero'><div class='kicker'>Paper-to-Course</div><h1>{title}</h1><div class='lead'>{lead}</div><div class='meta'>{pills}</div></header>{content}<div class='footer'>{footer}</div></div><script src='main.js'></script></body></html>"""


def esc(text: str) -> str:
    return (str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;"))


def write(path: Path, text: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding='utf-8')


def pills(course):
    return ''.join(f"<span class='pill'>{esc(x)}</span>" for x in [course['subtitle'], course['domain'], course['doi']])


def cards(items):
    return "<div class='grid cards'>" + ''.join(f"<article class='card'><h3>{esc(t)}</h3><p>{esc(d)}</p></article>" for t, d in items) + "</div>"


def timeline(items):
    body = ["<div class='timeline'>"]
    for year, title, desc in items:
        body.append(f"<div class='timeline-item'><div class='timeline-year'>{year}</div><strong>{esc(title)}</strong><div>{esc(desc)}</div></div>")
    body.append("</div>")
    return ''.join(body)


def table(headers, rows):
    head = ''.join(f"<th>{esc(h)}</th>" for h in headers)
    body = ''.join("<tr>" + ''.join(f"<td>{esc(cell)}</td>" for cell in row) + "</tr>" for row in rows)
    return f"<div class='table-wrap'><table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table></div>"


def bullets(items):
    return ''.join(f"<li>{esc(x)}</li>" for x in items)


def quiz_box(quiz):
    opts = ''.join(f"<button data-choice='{esc(v)}'>{esc(t)}</button>" for t, v in quiz['options'])
    return f"<div class='card quiz' data-quiz data-answer='{esc(quiz['answer'])}'><h3>Quick check</h3><p>{esc(quiz['question'])}</p>{opts}</div>"


def module_html(course, i):
    if i == 1:
        return f"<section id='module-1' class='module'><h2>Module 1 · 问题与动机</h2><p>{esc(course['summary'])}</p>{cards(course['cards'])}{quiz_box(course['quiz'])}</section>"
    if i == 2:
        return f"<section id='module-2' class='module'><h2>Module 2 · 发展脉络</h2>{timeline(course['timeline'])}</section>"
    if i == 3:
        return f"<section id='module-3' class='module'><h2>Module 3 · 主流方法对比</h2>{table(['Method', 'Strengths', 'Weaknesses'], course['compare'])}</section>"
    if i == 4:
        return f"<section id='module-4' class='module'><h2>Module 4 · 本文方法详解</h2><div class='grid two-col'><div class='card'><h3>Method flow</h3><ol>{''.join(f'<li><strong>{esc(num)}. {esc(title)}</strong> {esc(desc)}</li>' for num, title, desc in course['steps'])}</ol></div><div class='card'><h3>Key formula</h3><div class='quote'>{esc(course['formula'])}</div><ul>{''.join(f'<li><strong>{esc(a)}</strong> · {esc(b)}</li>' for a,b in course['formula_lines'])}</ul></div></div></section>"
    if i == 5:
        stats_html = "<div class='grid cards'>" + ''.join(f"<article class='card'><div style='font-size:28px;font-weight:800;color:var(--accent);'>{esc(v)}</div><h3>{esc(l)}</h3></article>" for v, l in course['stats']) + "</div>"
        bars_html = "<div class='grid cards'>" + ''.join(f"<div class='card'><strong>{esc(label)}</strong><div style='margin-top:10px;background:#f1ebe2;border-radius:999px;overflow:hidden;height:14px;'><div style='width:{val}%;height:100%;background:linear-gradient(90deg,var(--accent),#f27e57);'></div></div><small>{esc(note)} · {val}</small></div>" for label, val, note in course['bars']) + "</div>"
        return f"<section id='module-5' class='module'><h2>Module 5 · 实验与验证</h2>{stats_html}{bars_html}</section>"
    return f"<section id='module-6' class='module'><h2>Module 6 · 局限与展望</h2><div class='grid two-col'><div class='card'><h3>Current limitations</h3><ul>{bullets(course['limitations'])}</ul></div><div class='card'><h3>Future work</h3><ul>{bullets(course['future'])}</ul></div></div></section>"


def course_html(course):
    nav = ''.join(f"<a href='#module-{i}'>Module {i}</a>" for i in range(1, 7))
    content = ''.join(module_html(course, i) for i in range(1, 7))
    footer = f"<div class='quote'><strong>Source:</strong> {esc(course['title'])} | {esc(course['authors'])} | <a href='{esc(course['url'])}'>{esc(course['url'])}</a></div>"
    return BASE_HTML.format(title=esc(course['title']), nav=nav, lead=esc(course['summary']), pills=pills(course), content=content, footer=footer)


def readme(course):
    return textwrap.dedent(f'''
    # {course['title']}

    {course['subtitle']}

    ## Source
    - DOI: [{course['doi']}](https://doi.org/{course['doi']})
    - Article page: {course['url']}
    - Authors: {course['authors']}

    ## Course modules
    1. Problem and motivation
    2. Development timeline
    3. Method comparison
    4. Method details
    5. Experiments and validation
    6. Limitations and outlook

    ## Key takeaway
    {course['summary']}
    ''').strip() + '\n'


def slides_config(course):
    return {
        'title': course['title'],
        'subtitle': course['subtitle'],
        'source': course['url'],
        'slides': [
            {'type': 'title', 'title': course['title'], 'subtitle': course['subtitle'], 'note': course['source_note']},
            {'type': 'outline', 'title': '汇报提纲', 'items': ['问题与动机', '发展脉络', '方法对比', '本文方法', '实验结果', '局限与展望']},
            {'type': 'content', 'title': '研究背景', 'cards': [{'title': t, 'text': d} for t, d in course['cards']]},
            {'type': 'timeline', 'title': '发展脉络', 'items': [{'year': str(y), 'title': t, 'desc': d} for y, t, d in course['timeline']]},
            {'type': 'table', 'title': 'Benchmark 对比', 'headers': ['Method', 'Strengths', 'Weaknesses'], 'rows': course['compare']},
            {'type': 'flow', 'title': '方法流程', 'steps': [{'num': n, 'title': t, 'desc': d} for n, t, d in course['steps']]},
            {'type': 'formula', 'title': '核心公式与机理', 'formula': course['formula'], 'lines': [{'symbol': s, 'english': e} for s, e in course['formula_lines']]},
            {'type': 'stats', 'title': '核心结果', 'stats': [{'value': v, 'label': l} for v, l in course['stats']]},
            {'type': 'bars', 'title': '相对表现', 'items': [{'label': label, 'value': value, 'drop': note} for label, value, note in course['bars']]},
            {'type': 'content', 'title': '实验与验证', 'cards': [
                {'title': 'Method validation', 'text': 'Structure, thermal and cooling behavior are all checked.'},
                {'title': 'Indoor tests', 'text': 'Simulated skin and wearable scenarios are used.'},
                {'title': 'Outdoor tests', 'text': 'Sunlight or environment are incorporated when relevant.'},
                {'title': 'Stability', 'text': 'Cycle, wetting or degradation are quantified.'},
            ]},
            {'type': 'content', 'title': 'Mechanism', 'cards': [
                {'title': 'Heat transfer', 'text': 'The filler network changes phonon pathways and heat spreading.'},
                {'title': 'Interface control', 'text': 'Surface chemistry decides whether the network remains useful.'},
            ]},
            {'type': 'table', 'title': 'Key numbers to remember', 'headers': ['Metric', 'Value', 'Meaning'], 'rows': [
                ['Main thermal metric', course['stats'][0][0], course['stats'][0][1]],
                ['Secondary metric', course['stats'][1][0], course['stats'][1][1]],
                ['Cooling metric', course['stats'][2][0], course['stats'][2][1]],
                ['Durability / window', course['stats'][3][0], course['stats'][3][1]],
            ]},
            {'type': 'content', 'title': 'Application scenarios', 'cards': [
                {'title': 'Skin contact', 'text': 'Wearable patches and body surface cooling.'},
                {'title': 'Bioelectronics', 'text': 'Heat removal from LEDs or soft devices.'},
                {'title': 'Textiles / films', 'text': 'Large-area passive heat management.'},
                {'title': 'Reusable settings', 'text': 'Washing, cycling or reuse depending on the paper.'},
            ]},
            {'type': 'limitations', 'title': '局限与展望', 'limitations': course['limitations'], 'futureWork': course['future']},
            {'type': 'summary', 'title': '总结', 'items': [
                {'text': course['summary']},
                {'text': 'The main design principle is to align filler geometry, interfacial chemistry and matrix choice with the target cooling mode.'},
                {'text': f"Source: {course['url']}"},
            ]},
            {'type': 'content', 'title': 'References', 'cards': [{'title': 'DOI', 'text': course['doi']}, {'title': 'Article page', 'text': course['url']}]},
        ]
    }


def pptx_build(course, config, out):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_bar(slide, title, subtitle=None):
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
        band.fill.solid(); band.fill.fore_color.rgb = RGBColor(54, 69, 79); band.line.fill.background()
        box = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(12.2), Inches(0.5)).text_frame
        p = box.paragraphs[0]; r = p.add_run(); r.text = title; r.font.name = 'Arial'; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(255,255,255)
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.92), Inches(12), Inches(0.35)).text_frame
            sub.text = subtitle

    def body_box(slide, x, y, w, h, title, body):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 250, 244); shape.line.color.rgb = RGBColor(234, 223, 206)
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = title; r.font.name = 'Arial'; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = RGBColor(217,79,48)
        p = tf.add_paragraph(); p.text = body; p.font.name = 'Arial'; p.font.size = Pt(14)

    for slide_cfg in config['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = slide_cfg['type']
        if t == 'title':
            title_bar(slide, slide_cfg['title'], slide_cfg.get('subtitle'))
            body_box(slide, 0.7, 1.6, 11.9, 2.4, 'Course overview', course['summary'])
            body_box(slide, 0.7, 4.2, 11.9, 1.2, 'Source', course['url'])
        elif t == 'outline':
            title_bar(slide, slide_cfg['title'])
            body_box(slide, 0.5, 1.45, 12.2, 5.4, 'Outline', '\n'.join(f"{i+1}. {x}" for i, x in enumerate(slide_cfg['items'])))
        elif t == 'content':
            title_bar(slide, slide_cfg['title'])
            cards = slide_cfg['cards']
            for i, card in enumerate(cards):
                x = 0.55 + (i % 2) * 6.25
                y = 1.45 + (i // 2) * 2.3
                body_box(slide, x, y, 5.95, 1.95, card['title'], card['text'])
        elif t == 'timeline':
            title_bar(slide, slide_cfg['title'])
            body_box(slide, 0.5, 1.45, 12.2, 5.4, 'Timeline', '\n'.join(f"{x['year']} · {x['title']} — {x['desc']}" for x in slide_cfg['items']))
        elif t == 'table':
            title_bar(slide, slide_cfg['title'])
            rows = slide_cfg['rows']
            table_shape = slide.shapes.add_table(len(rows) + 1, len(slide_cfg['headers']), Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.8)).table
            for c, h in enumerate(slide_cfg['headers']):
                table_shape.cell(0, c).text = h
            for r, row in enumerate(rows, start=1):
                for c, val in enumerate(row):
                    table_shape.cell(r, c).text = str(val)
        elif t == 'flow':
            title_bar(slide, slide_cfg['title'])
            for i, step in enumerate(slide_cfg['steps']):
                body_box(slide, 0.55, 1.45 + i * 0.95, 12.1, 0.78, f"{step['num']}. {step['title']}", step['desc'])
        elif t == 'formula':
            title_bar(slide, slide_cfg['title'])
            body_box(slide, 0.5, 1.45, 12.2, 1.35, 'Formula', slide_cfg['formula'])
            body_box(slide, 0.5, 3.0, 12.2, 2.9, 'Interpretation', '\n'.join(f"{x['symbol']} — {x['english']}" for x in slide_cfg['lines']))
        elif t == 'stats':
            title_bar(slide, slide_cfg['title'])
            for i, st in enumerate(slide_cfg['stats']):
                body_box(slide, 0.55 + (i % 2) * 6.25, 1.5 + (i // 2) * 2.0, 5.95, 1.6, st['value'], st['label'])
        elif t == 'bars':
            title_bar(slide, slide_cfg['title'])
            maxv = max(item['value'] for item in slide_cfg['items'])
            for i, item in enumerate(slide_cfg['items']):
                y = 1.55 + i * 0.9
                label = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(2.2), Inches(0.3)).text_frame
                label.text = item['label']
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.7), Inches(y), Inches(8.7 * item['value'] / maxv), Inches(0.28))
                bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(217,79,48); bar.line.fill.background()
        elif t == 'limitations':
            title_bar(slide, slide_cfg['title'])
            body_box(slide, 0.5, 1.45, 6.1, 5.4, 'Current limitations', '\n'.join(f'- {x}' for x in slide_cfg['limitations']))
            body_box(slide, 6.75, 1.45, 6.1, 5.4, 'Future work', '\n'.join(f'- {x}' for x in slide_cfg['futureWork']))
        elif t == 'summary':
            title_bar(slide, slide_cfg['title'])
            body_box(slide, 0.5, 1.45, 12.2, 5.4, 'Summary', '\n'.join(x['text'] for x in slide_cfg['items']))

    prs.save(str(out))


def main():
    ROOT.mkdir(parents=True, exist_ok=True)
    manifest = []
    for course in COURSES:
        course_dir = ROOT / course['name']
        modules_dir = course_dir / 'modules'
        modules_dir.mkdir(parents=True, exist_ok=True)
        write(course_dir / 'styles.css', CSS)
        write(course_dir / 'main.js', JS)
        write(course_dir / '_base.html', '<!-- base template retained for compatibility -->')
        write(course_dir / '_footer.html', '<!-- footer template retained for compatibility -->')
        write(course_dir / 'build.sh', '#!/usr/bin/env bash\necho "Rebuild by rerunning g:/pythonProject/tools/generate_paper_courses.py"\n')
        write(course_dir / 'README.md', readme(course))
        cfg = slides_config(course)
        write(course_dir / 'slides-config.json', json.dumps(cfg, ensure_ascii=False, indent=2))
        write(course_dir / 'index.html', course_html(course))
        for i in range(1, 7):
            write(modules_dir / f'module-{i:02d}.html', module_html(course, i))
        pptx_build(course, cfg, course_dir / 'slides.pptx')
        manifest.append({'course': course['title'], 'path': str(course_dir), 'source': course['url']})

    write(ROOT / 'MANIFEST.md', '# Paper-to-Course Deliverables\n\n' + '\n'.join(f"- {x['course']} — {x['path']}" for x in manifest) + '\n')
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()