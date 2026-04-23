from __future__ import annotations

import json
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(r"g:/pythonProject/deliverables/paper-to-course-cn")

COURSES = [
    {
        "name": "2024_bnns_sf_cooling_patch_cn",
        "title": "高导热与辐射制冷协同设计的氮化硼纳米片/丝素蛋白薄膜用于个体热管理",
        "subtitle": "ACS Applied Materials & Interfaces 2024",
        "authors": "Juncheng Xia 等",
        "doi": "10.1021/acsami.3c16602",
        "url": "https://pubs.acs.org/doi/10.1021/acsami.3c16602",
        "summary": "该工作将BNNS与丝素蛋白耦合构建定向导热网络，在保证可穿戴与生物相容性的同时，实现热扩散与辐射制冷协同。",
        "domain": "个体热管理 / 辐射制冷 / 生物基导热薄膜",
        "timeline": [
            (2012, "丝素热输运研究兴起", "丝素材料从力学扩展到热功能研究。"),
            (2017, "BNNS复合导热膜发展", "二维片层定向导热成为核心路线。"),
            (2020, "可穿戴被动制冷加速", "服饰与薄膜场景并行推进。"),
            (2024, "BNNS/SF协同体系", "导热、反射与生物相容性一体化。"),
        ],
        "compare": [
            ["纯丝素薄膜", "柔性好、生物友好", "导热能力不足"],
            ["石墨烯纺织体系", "导热表现较强", "皮肤贴附与户外制冷存在折中"],
            ["BNNS聚合物膜", "面内导热突出", "基体生物友好性一般"],
            ["BNNS/SF（本文）", "定向导热+可重复使用+相容性好", "高填充可能影响柔软度"],
        ],
        "cards": [
            ["研究意义", "制冷能耗高，皮肤侧被动散热是低能耗路线。"],
            ["关键瓶颈", "丝素本体相容性好但热输运能力不够。"],
            ["核心思路", "以BNNS构建定向热通道，以丝素提供柔性与可持续基体。"],
        ],
        "steps": [
            ["1", "制备BNNS", "由h-BN剥离得到片层填料。"],
            ["2", "与丝素复配", "溶液混合后进行过滤成膜。"],
            ["3", "形成定向网络", "层状堆叠提升面内导热效率。"],
            ["4", "性能验证", "测试导热、光谱、循环稳定与生物安全。"],
        ],
        "formula": "λ = α × ρ × Cp；导热增强率 TCE = (λ复合 - λ基体) / λ基体 × 100%",
        "formula_lines": [
            ["λ", "等效热导率"],
            ["α", "热扩散率"],
            ["ρ", "密度"],
            ["Cp", "比热容"],
        ],
        "stats": [
            ["27.58 W m−1 K−1", "面内热导率"],
            ["1.77 W m−1 K−1", "厚向热导率"],
            ["0.89", "太阳反射率"],
            ["0.84", "中红外发射率"],
        ],
        "limitations": [
            "高填料比例可能带来柔性下降。",
            "大面积一致性制备仍需工程化优化。",
            "长期穿戴场景下耐磨与耐汗数据仍需扩展。",
        ],
        "future": [
            "优化界面耦合降低声子散射。",
            "发展可规模化定向成膜工艺。",
            "开展更长周期的真实服役测试。",
            "与纺织结构协同设计提升舒适性。",
        ],
    },
    {
        "name": "2020_silk_bn_nanofibers_cn",
        "title": "可调结构与性能的电纺丝素-氮化硼复合纳米纤维",
        "subtitle": "Polymers 2020",
        "authors": "Ye Xue, Xiao Hu",
        "doi": "10.3390/polym12051093",
        "url": "https://www.mdpi.com/2073-4360/12/5/1093",
        "summary": "该工作表明h-BN可显著调控电纺丝素纤维的形貌、二级结构与热稳定性，为生物基热功能纤维提供可调设计路径。",
        "domain": "电纺复合纤维 / 丝素蛋白 / h-BN填料",
        "timeline": [
            (2003, "电纺技术成熟", "纤维形貌可控成为常规手段。"),
            (2010, "丝素生物材料拓展", "丝素在组织工程与器件中广泛应用。"),
            (2015, "BN导热填料兴起", "电绝缘高导热填料体系快速发展。"),
            (2020, "丝素-BN纳米纤维", "结构与热稳定调控机制被系统验证。"),
        ],
        "compare": [
            ["纯丝素纤维", "生物友好、表面平整", "热稳定与功能窗口有限"],
            ["低含量BN复合", "对齐与结构改善", "增强幅度有限"],
            ["高含量BN复合", "耦合增强更明显", "粗糙度与工艺窗口收窄"],
            ["本文体系", "结构可调、机理清晰", "直接导热定量对比较少"],
        ],
        "cards": [
            ["研究意义", "电纺丝素可构建高比表面积网络，但热稳定性仍有提升空间。"],
            ["关键瓶颈", "填料增强需兼顾纤维可纺性与结构完整性。"],
            ["核心思路", "通过h-BN比例调控实现形貌-结构-热性质联动优化。"],
        ],
        "steps": [
            ["1", "配制丝素纺丝液", "建立稳定的电纺前驱体系。"],
            ["2", "引入h-BN", "按比例混合形成复合溶液。"],
            ["3", "电纺成纤", "获得不同结构特征的复合纤维膜。"],
            ["4", "热与结构表征", "通过SEM/FTIR/TGA/TM-DSC分析机制。"],
        ],
        "formula": "通过Tw、Tg、Td等特征温度迁移来刻画BN对丝素热行为的调控。",
        "formula_lines": [
            ["Tw", "束缚水相关热事件"],
            ["Tg", "玻璃化转变特征"],
            ["Td", "热降解特征"],
            ["纤维直径", "形貌调控的直接指标"],
        ],
        "stats": [
            ["约800 nm", "纯丝素纤维典型直径"],
            ["约400-750 nm", "复合后直径范围"],
            ["约120 °C", "纯丝素起始Tg"],
            ["约96 °C", "40% BN时最低Tg"],
        ],
        "limitations": [
            "论文重点在结构与热稳定，导热实测对比相对有限。",
            "高填充可能引入表面粗糙与缺陷。",
            "生物应用验证仍偏前瞻。",
        ],
        "future": [
            "补充导热各向异性与热扩散定量评估。",
            "建立结构调控与生物响应的关联模型。",
            "结合后处理提升有序化程度。",
            "探索多填料协同构筑多功能纤维。",
        ],
    },
    {
        "name": "2025_cd_bnoh_silma_patch_cn",
        "title": "CDs@BNOH填充丝素水凝胶的高导热与生物相容体表降温贴片",
        "subtitle": "ACS Applied Materials & Interfaces 2025",
        "authors": "Shuming Liu 等",
        "doi": "10.1021/acsami.5c09260",
        "url": "https://pubs.acs.org/doi/10.1021/acsami.5c09260",
        "summary": "该工作通过碳点修饰羟基化BN并引入丝素水凝胶，实现导热、体表降温、可降解与皮肤友好性的平衡。",
        "domain": "体表散热 / 丝素水凝胶 / 碳点修饰BN",
        "timeline": [
            (2020, "可穿戴被动散热", "水凝胶与织物体系并行发展。"),
            (2022, "生物相容凝胶强化", "皮肤贴附场景成为关键目标。"),
            (2024, "BNNS/SF薄膜验证", "薄膜路线证明了协同散热可行性。"),
            (2025, "CDs@BNOH/SilMA贴片", "界面修饰+可降解贴片方案成熟。"),
        ],
        "compare": [
            ["商用含水贴", "获取方便", "依赖蒸发，回温快"],
            ["纯SilMA凝胶", "柔软、生物友好", "热导率偏低"],
            ["BNOH/SilMA", "导热提升", "界面优化不足"],
            ["CDs@BNOH/SilMA", "导热与相容性平衡更好", "长期穿戴数据仍需扩展"],
        ],
        "cards": [
            ["研究意义", "体表散热材料需同时满足导热、贴肤和可持续要求。"],
            ["关键瓶颈", "高导热填料与生物基基体之间通常存在界面失配。"],
            ["核心思路", "以碳点调控BN表面界面，再与丝素甲基丙烯化凝胶耦合。"],
        ],
        "steps": [
            ["1", "制备碳点", "通过固相热合成得到CQDs。"],
            ["2", "构建CDs@BNOH", "在BNOH表面实现碳点修饰。"],
            ["3", "制备SilMA", "丝素甲基丙烯化用于光固化网络。"],
            ["4", "形成贴片", "复合填料注入凝胶并固化得到可贴附器件。"],
        ],
        "formula": "适度碳点修饰可延长有效声子输运路径，过度修饰会破坏导热网络完整性。",
        "formula_lines": [
            ["fwhm", "拉曼线宽用于表征声子寿命变化"],
            ["D%", "凝胶降解率"],
            ["ΔT", "相对空白组降温效果"],
            ["λ", "贴片热导率"],
        ],
        "stats": [
            ["0.77 W m−1 K−1", "贴片热导率"],
            ["15.38%", "导热提升幅度"],
            ["6.6 °C", "相对对照降温提升"],
            ["10 °C+", "室内工况峰值降温能力"],
        ],
        "limitations": [
            "热导率仍低于部分高填充薄膜体系。",
            "环境参数与贴附状态会影响户外表现。",
            "长期动态穿戴测试仍需补充。",
        ],
        "future": [
            "补充汗液、弯折和运动状态下的稳定性测试。",
            "进一步优化碳点修饰比例与分布。",
            "扩展到可穿戴电子散热应用场景。",
            "探索更多可降解基体与结构设计。",
        ],
    },
]

CSS = ":root{--bg:#faf7f2;--panel:#fffaf4;--text:#2d2620;--muted:#6b6057;--accent:#d94f30;--line:#eadfce;--shadow:0 10px 30px rgba(45,38,32,.08)}*{box-sizing:border-box}body{margin:0;font-family:Arial,'Microsoft YaHei',sans-serif;background:linear-gradient(180deg,#fcfbf8 0%,var(--bg) 100%);color:var(--text);line-height:1.6}.wrap{max-width:1160px;margin:0 auto;padding:0 20px 40px}.hero{padding:40px 0 22px;border-bottom:1px solid var(--line)}.kicker{text-transform:uppercase;letter-spacing:.18em;color:var(--accent);font-size:12px;font-weight:700}h1{font-size:clamp(30px,4vw,52px);line-height:1.05;margin:10px 0 12px}.lead{max-width:860px;color:var(--muted);font-size:17px}.meta{display:flex;flex-wrap:wrap;gap:10px;margin-top:18px}.pill{background:rgba(217,79,48,.09);color:#8a4d36;border:1px solid rgba(217,79,48,.18);border-radius:999px;padding:7px 12px;font-size:13px}.nav{position:sticky;top:0;z-index:20;background:rgba(250,247,242,.9);backdrop-filter:blur(8px);border-bottom:1px solid var(--line)}.nav-inner{display:flex;flex-wrap:wrap;gap:10px;padding:12px 0}.nav a{padding:8px 12px;border-radius:10px;color:var(--text);font-size:14px;border:1px solid transparent;text-decoration:none}.nav a.active,.nav a:hover{background:#fff;border-color:var(--line);box-shadow:var(--shadow)}.module{background:rgba(255,255,255,.86);border:1px solid var(--line);border-radius:20px;box-shadow:var(--shadow);padding:24px;margin-top:18px}.module h2{margin:0 0 8px;font-size:26px}.module p{color:var(--muted)}.grid{display:grid;gap:18px}.cards{grid-template-columns:repeat(auto-fit,minmax(220px,1fr))}.two-col{grid-template-columns:1.1fr .9fr}.card{background:var(--panel);border:1px solid var(--line);border-radius:16px;padding:18px}.card h3{margin:0 0 8px;font-size:18px}.table-wrap{overflow-x:auto;border-radius:14px;border:1px solid var(--line)}table{width:100%;border-collapse:collapse;background:#fff}th,td{padding:12px 14px;border-bottom:1px solid #eee7de;text-align:left;vertical-align:top}th{background:#fcf6ef}.timeline{position:relative;margin:14px 0;padding-left:24px}.timeline:before{content:'';position:absolute;left:9px;top:4px;bottom:4px;width:2px;background:rgba(217,79,48,.22)}.timeline-item{position:relative;margin-bottom:16px}.timeline-item:before{content:'';position:absolute;left:-24px;top:6px;width:18px;height:18px;border-radius:50%;background:var(--accent)}.timeline-year{font-weight:800;color:var(--accent)}.quote{border-left:4px solid var(--accent);padding:12px 14px;background:#fff6f0;border-radius:0 12px 12px 0}.footer{margin-top:24px;padding:20px 0 12px;color:var(--muted);border-top:1px solid var(--line);font-size:14px}@media(max-width:900px){.two-col{grid-template-columns:1fr}}"
JS = "window.addEventListener('DOMContentLoaded',()=>{const links=[...document.querySelectorAll('.nav a')];const sections=links.map(link=>document.querySelector(link.getAttribute('href'))).filter(Boolean);const observer=new IntersectionObserver(entries=>{const visible=entries.filter(entry=>entry.isIntersecting).sort((a,b)=>b.intersectionRatio-a.intersectionRatio)[0];if(!visible)return;links.forEach(link=>link.classList.toggle('active',link.getAttribute('href')==='#'+visible.target.id))},{threshold:[.2,.4,.6]});sections.forEach(sec=>observer.observe(sec));});"
BASE_HTML = "<!doctype html><html lang='zh-CN'><head><meta charset='utf-8'/><meta name='viewport' content='width=device-width,initial-scale=1'/><title>{title}</title><link rel='stylesheet' href='styles.css'/></head><body><div class='nav'><div class='wrap nav-inner'>{nav}</div></div><div class='wrap'><header class='hero'><div class='kicker'>Paper-to-Course 中文版</div><h1>{title}</h1><div class='lead'>{lead}</div><div class='meta'>{pills}</div></header>{content}<div class='footer'>{footer}</div></div><script src='main.js'></script></body></html>"


def esc(s: str) -> str:
    return str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def write(path: Path, text: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def mod(course: dict, idx: int) -> str:
    if idx == 1:
        cards = "".join(f"<article class='card'><h3>{esc(t)}</h3><p>{esc(d)}</p></article>" for t, d in course["cards"])
        return f"<section id='module-1' class='module'><h2>Module 1 · 问题与动机</h2><p>{esc(course['summary'])}</p><div class='grid cards'>{cards}</div></section>"
    if idx == 2:
        items = "".join(f"<div class='timeline-item'><div class='timeline-year'>{y}</div><strong>{esc(t)}</strong><div>{esc(d)}</div></div>" for y, t, d in course["timeline"])
        return f"<section id='module-2' class='module'><h2>Module 2 · 发展脉络</h2><div class='timeline'>{items}</div></section>"
    if idx == 3:
        rows = "".join("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in row) + "</tr>" for row in course["compare"])
        return "<section id='module-3' class='module'><h2>Module 3 · 主流方法对比</h2><div class='table-wrap'><table><thead><tr><th>方法</th><th>优点</th><th>不足</th></tr></thead><tbody>" + rows + "</tbody></table></div></section>"
    if idx == 4:
        steps = "".join(f"<li><strong>{esc(n)}. {esc(t)}</strong> {esc(d)}</li>" for n, t, d in course["steps"])
        lines = "".join(f"<li><strong>{esc(a)}</strong> · {esc(b)}</li>" for a, b in course["formula_lines"])
        return f"<section id='module-4' class='module'><h2>Module 4 · 本文方法详解</h2><div class='grid two-col'><div class='card'><h3>方法流程</h3><ol>{steps}</ol></div><div class='card'><h3>核心公式</h3><div class='quote'>{esc(course['formula'])}</div><ul>{lines}</ul></div></div></section>"
    if idx == 5:
        stats = "".join(f"<article class='card'><div style='font-size:28px;font-weight:800;color:var(--accent);'>{esc(v)}</div><h3>{esc(l)}</h3></article>" for v, l in course["stats"])
        return f"<section id='module-5' class='module'><h2>Module 5 · 实验与验证</h2><div class='grid cards'>{stats}</div></section>"
    lim = "".join(f"<li>{esc(x)}</li>" for x in course["limitations"])
    fut = "".join(f"<li>{esc(x)}</li>" for x in course["future"])
    return f"<section id='module-6' class='module'><h2>Module 6 · 局限与展望</h2><div class='grid two-col'><div class='card'><h3>当前局限</h3><ul>{lim}</ul></div><div class='card'><h3>未来方向</h3><ul>{fut}</ul></div></div></section>"


def build_pptx(course: dict, out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_bar(slide, title, subtitle=""):
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
        band.fill.solid(); band.fill.fore_color.rgb = RGBColor(54, 69, 79); band.line.fill.background()
        t = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(12), Inches(0.5)).text_frame
        p = t.paragraphs[0]
        p.text = title
        p.runs[0].font.name = "Arial"
        p.runs[0].font.size = Pt(23)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(255, 255, 255)
        if subtitle:
            s = slide.shapes.add_textbox(Inches(0.55), Inches(0.92), Inches(12), Inches(0.35)).text_frame
            s.text = subtitle

    def box(slide, x, y, w, h, title, body):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(255, 250, 244)
        sh.line.color.rgb = RGBColor(234, 223, 206)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(18)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(217, 79, 48)
        p2 = tf.add_paragraph(); p2.text = body

    slides = [
        ("title", "标题页", "来源说明"),
        ("outline", "汇报提纲", "6个模块"),
        ("bg", "研究背景", "问题、瓶颈、思路"),
        ("tl", "发展脉络", "时间线"),
        ("cmp", "方法对比", "主流方案比较"),
        ("flow", "方法流程", "制备与验证"),
        ("formula", "核心公式", "符号解释"),
        ("stats", "关键结果", "定量指标"),
        ("exp", "实验补充", "稳定性与场景验证"),
        ("mech", "机理理解", "界面与输运"),
        ("num", "关键数字", "便于复述"),
        ("app", "应用场景", "体表散热与器件"),
        ("risk", "当前局限", "边界条件"),
        ("future", "未来工作", "后续研究"),
        ("sum", "总结", "核心结论"),
        ("ref", "参考来源", "DOI 与链接"),
    ]

    for i, (_, t, d) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_bar(slide, f"{i:02d}. {t}", course["subtitle"])
        if i == 1:
            box(slide, 0.7, 1.6, 11.9, 2.3, "课程概述", course["summary"])
            box(slide, 0.7, 4.1, 11.9, 1.4, "论文来源", f"{course['title']}\n{course['authors']}\n{course['url']}")
        elif i == 2:
            box(slide, 0.5, 1.45, 12.2, 5.4, "提纲", "1. 问题与动机\n2. 发展脉络\n3. 主流方法对比\n4. 本文方法详解\n5. 实验与验证\n6. 局限与展望")
        elif i == 3:
            txt = "\n".join([f"- {x[0]}：{x[1]}" for x in course["cards"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "背景要点", txt)
        elif i == 4:
            txt = "\n".join([f"{y}：{t}（{d}）" for y, t, d in course["timeline"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "时间线", txt)
        elif i == 5:
            txt = "\n".join([f"- {r[0]} | 优点：{r[1]} | 不足：{r[2]}" for r in course["compare"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "方法对比", txt)
        elif i == 6:
            txt = "\n".join([f"{n}. {t}：{d}" for n, t, d in course["steps"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "流程", txt)
        elif i == 7:
            txt = course["formula"] + "\n" + "\n".join([f"{a}：{b}" for a, b in course["formula_lines"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "公式与符号", txt)
        elif i == 8:
            txt = "\n".join([f"{v}：{l}" for v, l in course["stats"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "核心指标", txt)
        elif i == 9:
            box(slide, 0.5, 1.45, 12.2, 5.4, "实验与验证", "围绕热导、温度响应、循环稳定与应用工况进行系统验证。")
        elif i == 10:
            box(slide, 0.5, 1.45, 12.2, 5.4, "机理理解", "热输运性能由填料网络拓扑与界面耦合共同决定。")
        elif i == 11:
            txt = "\n".join([f"- {v}（{l}）" for v, l in course["stats"]])
            box(slide, 0.5, 1.45, 12.2, 5.4, "关键数字", txt)
        elif i == 12:
            box(slide, 0.5, 1.45, 12.2, 5.4, "应用场景", "体表降温、可穿戴热管理、生物电子散热与材料扩展应用。")
        elif i == 13:
            box(slide, 0.5, 1.45, 12.2, 5.4, "当前局限", "\n".join([f"- {x}" for x in course["limitations"]]))
        elif i == 14:
            box(slide, 0.5, 1.45, 12.2, 5.4, "未来工作", "\n".join([f"- {x}" for x in course["future"]]))
        elif i == 15:
            box(slide, 0.5, 1.45, 12.2, 5.4, "总结", course["summary"])
        else:
            box(slide, 0.5, 1.45, 12.2, 5.4, "参考来源", f"DOI: {course['doi']}\n{course['url']}")

    prs.save(str(out_path))


def main():
    ROOT.mkdir(parents=True, exist_ok=True)
    manifest = []
    for course in COURSES:
        d = ROOT / course["name"]
        modules = d / "modules"
        modules.mkdir(parents=True, exist_ok=True)

        nav = "".join([f"<a href='#module-{i}'>Module {i}</a>" for i in range(1, 7)])
        pills = "".join([f"<span class='pill'>{esc(x)}</span>" for x in [course["subtitle"], course["domain"], course["doi"]]])
        content = "".join([mod(course, i) for i in range(1, 7)])
        footer = f"来源：{esc(course['title'])} | {esc(course['authors'])} | <a href='{esc(course['url'])}'>{esc(course['url'])}</a>"

        write(d / "styles.css", CSS)
        write(d / "main.js", JS)
        write(d / "_base.html", "<!-- 中文课程模板头部占位 -->")
        write(d / "_footer.html", "<!-- 中文课程模板尾部占位 -->")
        write(d / "build.sh", "#!/usr/bin/env bash\necho \"中文课程已生成，如需重建请执行 python g:/pythonProject/tools/generate_paper_courses_cn.py\"\n")
        write(d / "index.html", BASE_HTML.format(title=esc(course["title"]), nav=nav, lead=esc(course["summary"]), pills=pills, content=content, footer=footer))

        for i in range(1, 7):
            write(modules / f"module-{i:02d}.html", mod(course, i))

        cfg = {
            "title": course["title"],
            "subtitle": course["subtitle"],
            "source": course["url"],
            "slides": [{"type": "title", "title": course["title"], "subtitle": course["subtitle"]}] + [{"type": "content", "title": f"第{i}页"} for i in range(2, 17)]
        }
        write(d / "slides-config.json", json.dumps(cfg, ensure_ascii=False, indent=2))

        readme = textwrap.dedent(f"""
        # {course['title']}

        {course['subtitle']}

        ## 来源信息
        - DOI: [{course['doi']}](https://doi.org/{course['doi']})
        - 论文页面: {course['url']}
        - 作者: {course['authors']}

        ## 课程模块
        1. 问题与动机
        2. 发展脉络
        3. 主流方法对比
        4. 本文方法详解
        5. 实验与验证
        6. 局限与展望

        ## 核心结论
        {course['summary']}
        """).strip() + "\n"
        write(d / "README.md", readme)

        build_pptx(course, d / "slides.pptx")

        manifest.append({"course": course["title"], "path": str(d), "source": course["url"]})

    write(ROOT / "MANIFEST.md", "# 中文 Paper-to-Course 交付清单\n\n" + "\n".join([f"- {x['course']} — {x['path']}" for x in manifest]) + "\n")
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
