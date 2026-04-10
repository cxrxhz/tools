import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # 辅助函数：设置标题样式
    def set_title(shape, text):
        shape.text = text
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Microsoft YaHei UI'
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "公式计算程序介绍"
    title.text_frame.paragraphs[0].runs[0].font.name = 'Microsoft YaHei UI'
    title.text_frame.paragraphs[0].runs[0].font.bold = True
    title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 102, 204)
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    
    subtitle.text = "自动化解析、多层级计算与精美报表导出\n\n项目演示汇报"
    for r in subtitle.text_frame.paragraphs[0].runs:
        r.font.name = 'Microsoft YaHei UI'
        r.font.color.rgb = RGBColor(100, 100, 100)
    for r in subtitle.text_frame.paragraphs[2].runs:
        r.font.name = 'Microsoft YaHei UI'
        r.font.size = Pt(18)

    # 2. Program Overview Slide
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_title(slide.shapes.title, "程序概述 (Overview)")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "本程序是一个基于 Python 和 Tkinter 开发的完全交互式的公式计算工具。"
    p1 = tf.add_paragraph()
    p1.text = "1. 提供用户友好的界面来输入复杂的数学公式。"
    p2 = tf.add_paragraph()
    p2.text = "2. 利用 SymPy 构建强大的数学引擎，自动识别变量并解析公式结构。"
    p3 = tf.add_paragraph()
    p3.text = "3. 支持将未知变量继续分解为子公式进行深度计算。"
    p4 = tf.add_paragraph()
    p4.text = "4. 计算完成后，一键输出LaTeX公式图元与HTML计算报告。"
    
    for p in tf.paragraphs:
        p.space_after = Pt(14)
        for r in p.runs:
            r.font.name = 'Microsoft YaHei UI'
            r.font.size = Pt(20)

    # 3. Core Features Slide
    slide = prs.slides.add_slide(str_to_layout(prs, 1))
    set_title(slide.shapes.title, "核心特性 (Core Features)")
    
    tf = slide.placeholders[1].text_frame
    p0 = tf.paragraphs[0]
    p0.text = "智能变量解析："
    for r in p0.runs: r.font.bold = True
    p0_sub = tf.add_paragraph()
    p0_sub.text = "无需提前定义变量，输入公式后系统将自动提取除内置数学函数（如sin, cos, exp等）外的所有未知变量。"
    p0_sub.level = 1
    
    p1 = tf.add_paragraph()
    p1.text = "子公式级联推导："
    p1.space_before = Pt(14)
    for r in p1.runs: r.font.bold = True
    p1_sub = tf.add_paragraph()
    p1_sub.text = "在变量赋值阶段，若某个变量未知，可继续输入该变量的推导子公式。系统支持无限层级的嵌套解析直至计算出最终结果。"
    p1_sub.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "高精度计算引擎："
    p2.space_before = Pt(14)
    for r in p2.runs: r.font.bold = True
    p2_sub = tf.add_paragraph()
    p2_sub.text = "采用SymPy符号计算系统，支持浮点高精度计算，防止中间舍入误差积累。"
    p2_sub.level = 1
    
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = 'Microsoft YaHei UI'

    # 4. Result Output Slide
    slide = prs.slides.add_slide(str_to_layout(prs, 1))
    set_title(slide.shapes.title, "数据持久化与成果输出 (Data Output)")
    
    tf = slide.placeholders[1].text_frame
    p0 = tf.paragraphs[0]
    p0.text = "计算结果不仅仅以弹窗展示，更在后台静默生成高品质报告数据："
    
    p1 = tf.add_paragraph()
    p1.text = "排版级公式生成：调用 Matplotlib 将任意复杂度的公式渲染为 LaTeX 美化版 PNG 图片。"
    p1.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "原始数据留存：自动汇总所有经过计算的变量列表并导出至纯文本文件，方便接入其他系统。"
    p2.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "全流程 HTML 汇总报告：生成包含所有层级公式图片及变量对照表的 HTML 页面，方便展示和复用此计算过程。"
    p3.level = 1
    
    for p in tf.paragraphs:
        p.space_before = Pt(10)
        for r in p.runs:
            r.font.name = 'Microsoft YaHei UI'
            r.font.size = Pt(20)

    # 5. Tech Stack Slide
    slide = prs.slides.add_slide(str_to_layout(prs, 1))
    set_title(slide.shapes.title, "技术栈与架构实现 (Tech Stack)")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "界面开发："
    p0_sub = tf.add_paragraph()
    p0_sub.text = "使用 Python 内置的 Tkinter 构建图形用户界面，通过 Toplevel 实现多窗口调度与生命周期控制。"
    p0_sub.level = 1
    
    p1 = tf.add_paragraph()
    p1.text = "数学与逻辑核心："
    p1_sub = tf.add_paragraph()
    p1_sub.text = "采用 SymPy 模块 (symbols, parse_expr, N) 实现字符串的抽象语法树转换及代数计算。"
    p1_sub.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "图像与渲染："
    p2_sub = tf.add_paragraph()
    p2_sub.text = "结合 Matplotlib.pyplot 与 PIL 组件，基于 ByteIO 完成内存中的 LaTeX 图像渲染与本地编码落盘。"
    p2_sub.level = 1

    for p in tf.paragraphs:
        if p.level == 0:
            for r in p.runs: r.font.bold = True
        p.space_after = Pt(5)
        for r in p.runs:
            r.font.name = 'Microsoft YaHei UI'
            r.font.size = Pt(20) if p.level == 0 else Pt(16)

    output_path = os.path.join("g:\\pythonProject\\通用（未分类）\\公式计算", "程序功能演示.pptx")
    prs.save(output_path)
    print(f"PPT successfully created at {output_path}")

def str_to_layout(prs, idx):
    return prs.slide_layouts[idx]

if __name__ == '__main__':
    create_presentation()
