from pptx import Presentation
from pptx.util import Inches, Pt


def main():
    # 载入模板 PPT 文件
    prs = Presentation("template.pptx")

    # 选择目标布局（这里使用布局索引 2，对应母版中“标题幻灯片”）
    layout_index = 2
    layout = prs.slide_layouts[layout_index]

    # 定义图片占位符名称列表（必须和母版中的图片占位符名称一致）
    img_placeholder_names = ["Pic_Cais", "Pic_SEM1", "Pic_SEM2", "Pic_SEM3"]

    # 从母版（布局对象）中获取图片占位符的坐标信息
    layout_img_map = {}
    for shape in layout.shapes:
        try:
            # 如果能访问 placeholder_format 则表示该形状为占位符
            _ = shape.placeholder_format
            if shape.name in img_placeholder_names:
                layout_img_map[shape.name] = (shape.left, shape.top, shape.width, shape.height)
        except Exception:
            continue

    # 检查是否全部找到
    for name in img_placeholder_names:
        if name not in layout_img_map:
            print(f"警告：在布局中未找到图片占位符: {name}")

    # 同时从母版中获取小标题占位符的坐标，要求其名称为 "SmallTitle"
    layout_smalltitle_coord = None
    for shape in layout.shapes:
        try:
            _ = shape.placeholder_format
            if shape.name == "SmallTitle":
                layout_smalltitle_coord = (shape.left, shape.top, shape.width, shape.height)
                break
        except Exception:
            continue
    if layout_smalltitle_coord is None:
        print("警告：在布局中未找到小标题占位符 'SmallTitle'。请检查母版设置！")

    # 定义全部五组数据，每组数据包括：第一个元素为小标题，其后依次为对应4张图片的文件路径
    slide_data = [
        [  # 第一组：例如 B3-4-1
            "B3-4-1",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\蔡司\b3\4-1-1000-s.jpg",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-1-1.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-1-2.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-1-3.tif",
        ],
        [  # 第二组：例如 B3-4-3
            "B3-4-3",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\蔡司\b3\4-3-1000-s.jpg",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-3-1.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-3-2.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-3-3.tif",
        ],
        [  # 第三组：例如 B3-4-4
            "B3-4-4",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\蔡司\b3\4-4-1000-s.jpg",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-4-1.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-4-2.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B3\4-4-3.tif",
        ],
        [  # 第四组：例如 B4-4-1
            "B4-4-1",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\蔡司\b4\4-1-1000-s.jpg",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-1-1.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-1-2.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-1-3.tif",
        ],
        [  # 第五组：例如 B4-4-3
            "B4-4-3",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\蔡司\b4\4-3-1000-s.jpg",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-3-1.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-3-2.tif",
            r"C:\Users\xhz\OneDrive - tongji.edu.cn\实验数据\SiNW\SEM\B4\4-3-3.tif",
        ],
    ]

    # 针对每组数据生成幻灯片
    for group in slide_data:
        small_title = group[0]  # 小标题文本
        images = group[1:]  # 4 张图片路径

        # 新增幻灯片，采用目标布局
        slide = prs.slides.add_slide(layout)

        # 如果在母版中找到了小标题占位符坐标，则直接添加文本框到相同位置
        if layout_smalltitle_coord is not None:
            st_left, st_top, st_width, st_height = layout_smalltitle_coord
            tb = slide.shapes.add_textbox(st_left, st_top, st_width, st_height)
            tb.text_frame.text = small_title
            # 由于格式在母版中预先定义，故这里不再手动设置格式；你可以根据需要调整：
            for paragraph in tb.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(28)
        else:
            print("未能获得小标题的母版坐标，未生成小标题。")

        # 根据母版中图片占位符坐标添加图片
        for idx, img_ph_name in enumerate(img_placeholder_names):
            if img_ph_name in layout_img_map:
                left, top, width, height = layout_img_map[img_ph_name]
                slide.shapes.add_picture(images[idx], left, top, width, height)
            else:
                print(f"未在布局中找到图片占位符: {img_ph_name}")

    # 保存生成的 PPT 文件
    prs.save("generated_presentation_with_template_titles.pptx")
    print("所有幻灯片已生成！")


if __name__ == '__main__':
    main()
